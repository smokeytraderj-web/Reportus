"""Generate the branded Client Deck intermediate PPTX from validated data."""

from __future__ import annotations

import importlib.util
import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config.settings import PROJECT_ROOT


def _load_style():
    path = PROJECT_ROOT / "skills" / "client-deck-builder" / "style.py"
    spec = importlib.util.spec_from_file_location("reportus_client_deck_style", path)
    if spec is None or spec.loader is None:
        raise RuntimeError("Client Deck style resource is unavailable.")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


style = _load_style()


@dataclass(frozen=True, slots=True)
class AllocationRow:
    label: str
    value: float


@dataclass(frozen=True, slots=True)
class ClientDeckData:
    client_name: str
    period: str
    as_of: str
    allocation: tuple[AllocationRow, ...]
    risk_metrics: dict[str, str]
    sector_performance: dict[str, float]
    sector_portfolio: dict[str, float]
    sector_benchmark: dict[str, float]
    contributors: tuple[tuple[str, str, str, str], ...]
    detractors: tuple[tuple[str, str, str, str], ...]
    earnings_years: tuple[str, ...]
    earnings_values: tuple[float, ...]
    earnings_notes: tuple[str, ...]
    optional_sections: dict[str, tuple[str, ...]] = field(default_factory=dict)
    sources: dict[str, str] = field(default_factory=dict)


def _base_slide(prs: Presentation, title: str, page: int, total: int, as_of: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style.add_rect(slide, 0, 0, style.SLIDE_W, style.SLIDE_H, fill_rgb=style.WHITE)
    style.add_content_header(slide, title, page, total, meta_placeholder=as_of)
    return slide


def _source(slide, text: str) -> None:
    if text:
        style.add_text(slide, Inches(0.75), Inches(6.82), Inches(11.4), Inches(0.22),
                       f"Source: {text}", 8, style.BODY_GRAY, font=style.BODY_FONT, italic=True)


def _risk_metric(metrics: dict[str, str], *labels: str) -> str:
    normalized = {
        re.sub(r"[^a-z0-9%]", "", key.casefold()): value
        for key, value in metrics.items()
    }
    for label in labels:
        value = normalized.get(re.sub(r"[^a-z0-9%]", "", label.casefold()))
        if value:
            return value
    return ""


def _risk_number(text: str) -> int | None:
    match = re.search(r"\d{1,3}", text or "")
    if match is None:
        return None
    value = int(match.group())
    return value if 1 <= value <= 99 else None


def _add_riskalyze_snapshot(slide, data: ClientDeckData) -> None:
    """Recreate Riskalyze's information hierarchy in the GSWM visual system."""

    metrics = data.risk_metrics
    portfolio_total = _risk_metric(metrics, "Portfolio total")
    risk_number = _risk_number(_risk_metric(metrics, "Risk"))
    loss_amount = _risk_metric(metrics, "Historical loss", "95% historical loss")
    loss_percent = _risk_metric(metrics, "Historical loss %", "95% historical loss %")
    gain_amount = _risk_metric(metrics, "Historical gain", "95% historical gain")
    gain_percent = _risk_metric(metrics, "Historical gain %", "95% historical gain %")

    style.add_text(slide, Inches(.78), Inches(1.55), Inches(2.4), Inches(.2),
                   "PORTFOLIO TOTAL", 8, style.GOLD, font=style.BODY_FONT, bold=True, spacing=.7)
    style.add_text(slide, Inches(.78), Inches(1.82), Inches(4.2), Inches(.52),
                   portfolio_total, 26, style.NAVY, font=style.TITLE_FONT, bold=True)

    risk_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(11.35), Inches(1.48), Inches(1.05), Inches(1.05)
    )
    risk_circle.fill.background()
    risk_circle.line.color.rgb = style.GOLD
    risk_circle.line.width = Pt(1.6)
    style.no_shadow(risk_circle)
    style.add_text(slide, Inches(11.48), Inches(1.62), Inches(.8), Inches(.16),
                   "RISK", 8, style.GOLD, font=style.BODY_FONT, bold=True,
                   align=PP_ALIGN.CENTER, spacing=.5)
    style.add_text(slide, Inches(11.48), Inches(1.82), Inches(.8), Inches(.42),
                   str(risk_number or "—"), 25, style.NAVY, font=style.TITLE_FONT,
                   bold=True, align=PP_ALIGN.CENTER)

    total_value = sum(row.value for row in data.allocation)
    categories = [row.label for row in data.allocation]
    values = [row.value / total_value for row in data.allocation]
    style.add_donut_chart(
        slide, Inches(.72), Inches(2.38), Inches(3.15), Inches(3.72), categories, values,
        hole_size=68,
    )
    legend_top = 2.74
    for index, (label, value) in enumerate(zip(categories[:6], values[:6])):
        top = legend_top + index * .46
        style.add_rect(
            slide, Inches(3.5), Inches(top + .04), Inches(.10), Inches(.10),
            fill_rgb=style.CHART_COLORS[index % len(style.CHART_COLORS)],
        )
        style.add_text(slide, Inches(3.68), Inches(top), Inches(1.5), Inches(.18),
                       label, 8.5, style.NAVY, font=style.BODY_FONT)
        style.add_text(slide, Inches(4.66), Inches(top), Inches(.52), Inches(.18),
                       f"{value:.2%}", 8.5, style.NAVY, font=style.BODY_FONT,
                       bold=True, align=PP_ALIGN.RIGHT)

    style.add_text(slide, Inches(5.42), Inches(2.42), Inches(6.95), Inches(.2),
                   "95% HISTORICAL RANGE (6 MONTHS)", 8, style.GOLD,
                   font=style.BODY_FONT, bold=True, spacing=.65)
    if loss_amount and gain_amount:
        style.add_text(slide, Inches(5.42), Inches(2.74), Inches(2.85), Inches(.35),
                       loss_amount, 19, style.RED, font=style.TITLE_FONT, bold=True,
                       align=PP_ALIGN.CENTER)
        style.add_text(slide, Inches(9.45), Inches(2.74), Inches(2.85), Inches(.35),
                       gain_amount, 19, style.GREEN, font=style.TITLE_FONT, bold=True,
                       align=PP_ALIGN.CENTER)
        style.add_text(slide, Inches(5.42), Inches(3.10), Inches(2.85), Inches(.2),
                       loss_percent, 9, style.BODY_GRAY, font=style.BODY_FONT,
                       align=PP_ALIGN.CENTER)
        style.add_text(slide, Inches(9.45), Inches(3.10), Inches(2.85), Inches(.2),
                       gain_percent, 9, style.BODY_GRAY, font=style.BODY_FONT,
                       align=PP_ALIGN.CENTER)
        style.add_rect(slide, Inches(5.42), Inches(3.48), Inches(3.45), Pt(4), fill_rgb=style.RED)
        style.add_rect(slide, Inches(8.87), Inches(3.48), Inches(3.43), Pt(4), fill_rgb=style.GREEN)
        style.add_text(slide, Inches(5.42), Inches(3.60), Inches(1), Inches(.18),
                       "5%", 8, style.BODY_GRAY, font=style.BODY_FONT)
        style.add_text(slide, Inches(11.30), Inches(3.60), Inches(1), Inches(.18),
                       "95%", 8, style.BODY_GRAY, font=style.BODY_FONT, align=PP_ALIGN.RIGHT)

    detail_metrics = [
        ("Annual dividend", _risk_metric(metrics, "Annual dividend")),
        ("Max drawdown", _risk_metric(metrics, "Max drawdown")),
        ("Annual range midpoint", _risk_metric(metrics, "Annual range midpoint")),
        ("Expense ratio", _risk_metric(metrics, "Expense ratio", "Portfolio costs")),
    ]
    visible_metrics = [(label, value) for label, value in detail_metrics if value]
    if visible_metrics:
        card_width = 6.88 / len(visible_metrics)
        for index, (label, value) in enumerate(visible_metrics):
            left = 5.42 + index * card_width
            style.add_rect(
                slide, Inches(left), Inches(4.15), Inches(card_width - .14), Inches(1.18),
                fill_rgb=style.LIGHT_ROW,
            )
            style.add_text(slide, Inches(left + .13), Inches(4.33), Inches(card_width - .4), Inches(.28),
                           label.upper(), 7.5, style.GOLD, font=style.BODY_FONT,
                           bold=True, spacing=.45)
            style.add_text(slide, Inches(left + .13), Inches(4.70), Inches(card_width - .4), Inches(.36),
                           value, 18, style.NAVY, font=style.TITLE_FONT, bold=True)


def build_client_deck(data: ClientDeckData, output_path: Path) -> Path:
    """Build an editable branded deck; PDF conversion is a separate verified step."""

    optional_titles = {
        "rmd": "RMD Report",
        "529": "Portfolio Summary - 529 Accounts",
        "annuity": "Annuity Review",
    }
    sections = [
        "Overall Asset Allocation", "Risk Snapshot", "Sector YTD Performance",
        "Equity Sector Exposure", "Attribution Report", "S&P 500 Earnings Expectations",
    ] + [optional_titles[key] for key in ("rmd", "529", "annuity") if key in data.optional_sections]
    total = len(sections) + 2
    prs = Presentation()
    prs.slide_width, prs.slide_height = style.SLIDE_W, style.SLIDE_H
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    style.add_gradient_bg(cover)
    style.add_crest(cover, style.SLIDE_W / 2, Inches(0.75))
    style.add_text(cover, Inches(1), Inches(3.55), Inches(9.5), Inches(.3),
                   "PRIVATE CLIENT REVIEW", 11, style.GOLD, bold=True, spacing=1.5)
    style.add_text(cover, Inches(1), Inches(3.9), Inches(10.5), Inches(.85),
                   data.client_name, 40, style.WHITE, font=style.TITLE_FONT, bold=True)
    style.add_text(cover, Inches(1), Inches(4.62), Inches(10), Inches(.5),
                   f"Portfolio Review · {data.period}", 22, style.WHITE, font=style.TITLE_FONT)
    style.add_gold_rule(cover, Inches(1), Inches(5.28), Inches(1.6))
    style.add_toc_slide(prs, blank, [(index + 3, title) for index, title in enumerate(sections)], total)

    page = 3
    slide = _base_slide(prs, sections[0], page, total, data.as_of)
    total_value = sum(row.value for row in data.allocation)
    categories = [row.label for row in data.allocation]
    values = [row.value / total_value for row in data.allocation]
    style.add_donut_chart(slide, Inches(.75), Inches(1.65), Inches(4.2), Inches(4.7), categories, values)
    rows = [[row.label, f"${row.value:,.2f}", f"{row.value / total_value:.2%}"] for row in data.allocation]
    shape = slide.shapes.add_table(len(rows) + 1, 3, Inches(5.2), Inches(1.75), Inches(7.3), Inches(3.8))
    style.style_table(shape.table, ["Asset Class", "Value", "%"], rows, [.55, .27, .18], Inches(7.3), align_from_col=1, header_align_from_col=1)
    _source(slide, data.sources.get("allocation", ""))

    page += 1
    slide = _base_slide(prs, sections[1], page, total, data.as_of)
    _add_riskalyze_snapshot(slide, data)
    _source(slide, data.sources.get("risk", ""))

    page += 1
    slide = _base_slide(prs, sections[2], page, total, data.as_of)
    style.add_bar_chart(slide, Inches(.75), Inches(1.65), Inches(11.8), Inches(4.9), list(data.sector_performance), list(data.sector_performance.values()))
    _source(slide, data.sources.get("sector_performance", ""))

    page += 1
    slide = _base_slide(prs, sections[3], page, total, data.as_of)
    cats = list(data.sector_portfolio)
    portfolio = [data.sector_portfolio[key] for key in cats]
    benchmark = [data.sector_benchmark[key] for key in cats]
    style.add_light_comparison_table(slide, Inches(.75), Inches(1.55), Inches(11.8), Inches(.55), Inches(.5), Inches(2.0), cats, [data.client_name, "S&P 500"], [portfolio, benchmark])
    style.add_diff_bars_manual(slide, Inches(.75), Inches(3.35), Inches(11.8), Inches(2.8), cats, [a - b for a, b in zip(portfolio, benchmark)])
    _source(slide, data.sources.get("sector_exposure", ""))

    page += 1
    slide = _base_slide(prs, sections[4], page, total, data.as_of)
    for left, title, rows in ((.75, "Top Contributors", data.contributors), (6.75, "Top Detractors", data.detractors)):
        style.add_text(slide, Inches(left), Inches(1.55), Inches(5.6), Inches(.35), title, 17, style.NAVY, font=style.TITLE_FONT, bold=True)
        table = slide.shapes.add_table(len(rows) + 1, 4, Inches(left), Inches(2.0), Inches(5.8), Inches(3.9)).table
        style.style_table(table, ["Symbol", "Holding", "Return", "Contrib."], rows, [.16, .48, .18, .18], Inches(5.8), header_size=9, body_size=8.5, align_from_col=2, header_align_from_col=2)
    _source(slide, data.sources.get("attribution", ""))

    page += 1
    slide = _base_slide(prs, sections[5], page, total, data.as_of)
    style.add_bullets(slide, Inches(.8), Inches(1.75), Inches(3.5), Inches(4.5), data.earnings_notes, size=14)
    style.add_axis_bar_chart(slide, Inches(4.55), Inches(1.7), Inches(7.6), Inches(4.8), data.earnings_years, data.earnings_values, style.CHART_COLORS)
    _source(slide, data.sources.get("earnings", ""))

    for key in ("rmd", "529", "annuity"):
        if key not in data.optional_sections:
            continue
        page += 1
        slide = _base_slide(prs, optional_titles[key], page, total, data.as_of)
        style.add_bullets(slide, Inches(.9), Inches(1.7), Inches(11.4), Inches(4.8), data.optional_sections[key], size=15, space_after=14)
        _source(slide, data.sources.get(key, ""))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
