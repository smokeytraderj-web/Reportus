"""Build portfolio workbooks as first-class GSWM client decks."""

from __future__ import annotations

from collections import OrderedDict
from pathlib import Path
from typing import TYPE_CHECKING
from urllib.parse import urlsplit

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches

from generators.client_deck import style

if TYPE_CHECKING:
    from generators.portfolio_pdf import PortfolioHolding, PortfolioWorkbookData


def _compact(text: object, limit: int) -> str:
    value = " ".join(str(text or "").split())
    if len(value) <= limit:
        return value
    return value[: max(1, limit - 1)].rstrip(" ,.;:-") + "…"


def _assumption(data: "PortfolioWorkbookData", label: str, fallback: str = "Not supplied") -> str:
    wanted = "".join(character for character in label.casefold() if character.isalnum())
    for key, value in data.assumptions:
        normalized = "".join(character for character in key.casefold() if character.isalnum())
        if normalized == wanted:
            return str(value)
    return fallback


def _as_of(data: "PortfolioWorkbookData") -> str:
    timestamp = _assumption(data, "Price timestamp", "")
    if timestamp:
        concise = timestamp.split(";", 1)[0].strip()
        return concise if concise.casefold().startswith("as of") else f"As of {concise}"
    for part in data.subtitle.split("|"):
        if "as of" in part.casefold():
            return part.strip()
    return "As of date supplied in workbook"


def _money(value: float) -> str:
    return f"${value:,.2f}"


def _percent(value: float, decimals: int = 1) -> str:
    return f"{value:.{decimals}%}"


def _base_slide(prs: Presentation, title: str, page: int, total: int, as_of: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style.add_rect(slide, 0, 0, style.SLIDE_W, style.SLIDE_H, fill_rgb=style.WHITE)
    style.add_content_header(slide, title, page, total, meta_placeholder=as_of)
    return slide


def _source(slide, source_label: str, detail: str = "") -> None:
    value = source_label.strip()
    if detail:
        value = f"{value} · {detail}" if value else detail
    if value:
        style.add_text(
            slide, Inches(0.75), Inches(6.82), Inches(11.4), Inches(0.22),
            f"Source: {value}", 8, style.BODY_GRAY, font=style.BODY_FONT, italic=True,
        )


def _table(slide, headers, rows, left, top, width, height, weights, *,
           header_size=10, body_size=9, align_from_col=None, header_align_from_col=None):
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height)
    )
    style.style_table(
        shape.table, headers, rows, weights, Inches(width),
        header_size=header_size, body_size=body_size,
        align_from_col=align_from_col, header_align_from_col=header_align_from_col,
    )
    return shape.table


def _sleeves(holdings: tuple["PortfolioHolding", ...]):
    result: OrderedDict[str, dict[str, float]] = OrderedDict()
    for item in holdings:
        bucket = result.setdefault(item.sleeve or "Other", {"target": 0.0, "now": 0.0, "monthly": 0.0})
        bucket["target"] += item.target
        bucket["now"] += item.invest_now
        bucket["monthly"] += item.monthly
    return result


def _kpi(slide, left: float, label: str, value: str, detail: str = "") -> None:
    style.add_rect(slide, Inches(left), Inches(1.65), Inches(2.72), Inches(1.05), fill_rgb=style.LIGHT_ROW)
    style.add_text(slide, Inches(left + .16), Inches(1.79), Inches(2.4), Inches(.2),
                   label.upper(), 8, style.GOLD, font=style.BODY_FONT, bold=True, spacing=.7)
    style.add_text(slide, Inches(left + .16), Inches(2.05), Inches(2.4), Inches(.37),
                   value, 21, style.NAVY, font=style.TITLE_FONT, bold=True)
    if detail:
        style.add_text(slide, Inches(left + .16), Inches(2.43), Inches(2.4), Inches(.16),
                       detail, 7.5, style.BODY_GRAY, font=style.BODY_FONT, italic=True)


def build_portfolio_client_deck(
    data: "PortfolioWorkbookData",
    output_path: Path,
    *,
    client_name: str,
    report_title: str,
    period_label: str,
    source_label: str,
) -> Path:
    """Render structured workbook content through the official client-deck visual system."""

    rationale_chunks = [data.holdings[index:index + 6] for index in range(0, len(data.holdings), 6)]
    research_chunks = [data.research_updates[index:index + 4] for index in range(0, len(data.research_updates), 4)]
    source_chunks = [data.sources[index:index + 10] for index in range(0, len(data.sources), 10)]

    sections = ["Executive Summary", "Target Allocation", "Funding Plan"]
    sections.extend(
        "Investment Rationale" if index == 0 else "Investment Rationale (continued)"
        for index in range(len(rationale_chunks))
    )
    sections.append("Risk & Implementation")
    sections.extend(
        "Recent Research Checks" if index == 0 else "Recent Research Checks (continued)"
        for index in range(len(research_chunks))
    )
    sections.extend(
        "Sources & Data Audit" if index == 0 else "Sources & Data Audit (continued)"
        for index in range(len(source_chunks))
    )
    total = len(sections) + 2
    as_of = _as_of(data)

    prs = Presentation()
    prs.slide_width, prs.slide_height = style.SLIDE_W, style.SLIDE_H
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    style.add_gradient_bg(cover)
    style.add_crest(cover, style.SLIDE_W / 2, Inches(.75))
    style.add_text(cover, Inches(1), Inches(3.48), Inches(10), Inches(.3),
                   "PRIVATE CLIENT REVIEW", 11, style.GOLD, bold=True, spacing=1.5)
    style.add_text(cover, Inches(1), Inches(3.84), Inches(10.8), Inches(.76),
                   client_name, 40, style.WHITE, font=style.TITLE_FONT, bold=True)
    style.add_text(cover, Inches(1), Inches(4.56), Inches(10.8), Inches(.52),
                   f"{report_title} · {period_label}", 22, style.WHITE, font=style.TITLE_FONT)
    style.add_gold_rule(cover, Inches(1), Inches(5.25), Inches(1.6))
    style.add_text(cover, Inches(1), Inches(5.48), Inches(10.8), Inches(.42),
                   _compact(data.title, 110), 12, style.MUTED_ON_NAVY, font=style.BODY_FONT, italic=True)

    style.add_toc_slide(prs, blank, [(number + 3, title) for number, title in enumerate(sections)], total)

    page = 3
    initial = sum(item.invest_now for item in data.holdings)
    monthly = sum(item.monthly for item in data.holdings)
    weighted_fee = sum(item.target * item.expense_ratio for item in data.holdings)
    slide = _base_slide(prs, "Executive Summary", page, total, as_of)
    for left, label, value, detail in (
        (.75, "Initial funding", _money(initial), "Allocated across all holdings"),
        (3.75, "Monthly funding", _money(monthly), "Rules-based contributions"),
        (6.75, "Target invested", _percent(sum(item.target for item in data.holdings), 0), "No unassigned allocation"),
        (9.75, "Weighted ETF fee", _percent(weighted_fee, 3), "Portfolio-weighted expense ratio"),
    ):
        _kpi(slide, left, label, value, detail)
    style.add_text(slide, Inches(.75), Inches(3.05), Inches(7.35), Inches(.34),
                   "Portfolio mandate", 18, style.NAVY, font=style.TITLE_FONT, bold=True)
    style.add_text(slide, Inches(.75), Inches(3.48), Inches(7.35), Inches(2.65),
                   _compact(data.thesis, 900), 12, style.DARK_TEXT, font=style.BODY_FONT, line_spacing=1.18)
    sleeve_rows = [[_compact(label, 28), _percent(values["target"])] for label, values in _sleeves(data.holdings).items()]
    style.add_text(slide, Inches(8.45), Inches(3.05), Inches(4.1), Inches(.34),
                   "Portfolio structure", 18, style.NAVY, font=style.TITLE_FONT, bold=True)
    _table(slide, ["Sleeve", "Target"], sleeve_rows, 8.45, 3.48, 4.1, 2.1, [.72, .28],
           header_size=9, body_size=8.8, align_from_col=1, header_align_from_col=1)
    style.add_text(slide, Inches(8.45), Inches(5.78), Inches(4.0), Inches(.45),
                   _compact(data.subtitle, 140), 9, style.BODY_GRAY, font=style.BODY_FONT, italic=True)
    _source(slide, source_label, "Portfolio and assumptions tabs")

    page += 1
    slide = _base_slide(prs, "Target Allocation", page, total, as_of)
    sleeves = _sleeves(data.holdings)
    categories = [_compact(label, 24) for label in sleeves]
    weights = [values["target"] for values in sleeves.values()]
    style.add_donut_chart(slide, Inches(.65), Inches(1.58), Inches(4.55), Inches(4.85), categories, weights)
    allocation_rows = [
        [_compact(label, 30), _percent(values["target"]), _money(values["now"]), _money(values["monthly"])]
        for label, values in sleeves.items()
    ]
    allocation_rows.append(["Total", _percent(sum(weights), 0), _money(initial), _money(monthly)])
    _table(slide, ["Sleeve", "Target", "Invest now", "Monthly"], allocation_rows,
           5.25, 1.72, 7.33, 3.2, [.43, .18, .2, .19], header_size=10, body_size=9.3,
           align_from_col=1, header_align_from_col=1)
    style.add_rect(slide, Inches(5.25), Inches(5.18), Inches(7.33), Inches(.92), fill_rgb=style.LIGHT_ROW)
    style.add_text(slide, Inches(5.47), Inches(5.36), Inches(6.9), Inches(.5),
                   "Allocation is organized by distinct return drivers; fixed targets direct new cash toward relative laggards.",
                   11, style.NAVY, font=style.BODY_FONT, italic=True, anchor=MSO_ANCHOR.MIDDLE)
    _source(slide, source_label, "Target allocation and contribution schedule")

    page += 1
    slide = _base_slide(prs, "Funding Plan", page, total, as_of)
    funding_rows = [
        [item.ticker, _compact(item.holding, 35), _percent(item.target), _money(item.invest_now),
         _money(item.monthly), _money(item.price)]
        for item in data.holdings
    ]
    funding_rows.append(["TOTAL", "", _percent(sum(item.target for item in data.holdings), 0),
                         _money(initial), _money(monthly), "—"])
    _table(slide, ["Ticker", "Holding", "Target", "Invest now", "Monthly", "Price"], funding_rows,
           .75, 1.58, 11.83, 4.95, [.09, .37, .12, .15, .14, .13], header_size=9.5, body_size=8.4,
           align_from_col=2, header_align_from_col=2)
    _source(slide, source_label, "Prices and funding amounts reproduced from workbook")

    for index, chunk in enumerate(rationale_chunks):
        page += 1
        title = "Investment Rationale" if index == 0 else "Investment Rationale (continued)"
        slide = _base_slide(prs, title, page, total, as_of)
        rows = [[item.ticker, _compact(item.holding, 29), _compact(item.valuation, 27), _compact(item.reason, 235)] for item in chunk]
        _table(slide, ["Ticker", "Holding", "Valuation", "Portfolio role and rationale"], rows,
               .75, 1.58, 11.83, 4.98, [.08, .24, .18, .50], header_size=9.5, body_size=8.2)
        _source(slide, source_label, "Valuation and rationale fields reproduced from workbook")

    page += 1
    slide = _base_slide(prs, "Risk & Implementation", page, total, as_of)
    for top, label, value in (
        (1.58, "BIGGEST RISK", data.biggest_risk),
        (3.05, "IMPLEMENTATION RULE", data.implementation_rule),
    ):
        style.add_rect(slide, Inches(.75), Inches(top), Inches(7.55), Inches(1.22), fill_rgb=style.LIGHT_ROW)
        style.add_text(slide, Inches(.98), Inches(top + .15), Inches(7.08), Inches(.2),
                       label, 8, style.GOLD, font=style.BODY_FONT, bold=True, spacing=.7)
        style.add_text(slide, Inches(.98), Inches(top + .42), Inches(7.08), Inches(.66),
                       _compact(value, 620), 10.5, style.DARK_TEXT, font=style.BODY_FONT, line_spacing=1.05)
    assumption_rows = [[_compact(label, 28), _compact(value, 120)] for label, value in data.assumptions[:7]]
    style.add_text(slide, Inches(8.62), Inches(1.58), Inches(3.95), Inches(.3),
                   "Key assumptions", 18, style.NAVY, font=style.TITLE_FONT, bold=True)
    _table(slide, ["Control", "Workbook value"], assumption_rows, 8.62, 2.02, 3.96, 3.95,
           [.38, .62], header_size=8.5, body_size=7.7)
    style.add_text(slide, Inches(.75), Inches(4.58), Inches(7.55), Inches(1.18),
                   "Monitoring standard\nReview no more than monthly. Reassess for thesis breaks, governance concerns, permanent competitive impairment, or fund methodology changes—not price volatility alone.",
                   10.5, style.NAVY, font=style.BODY_FONT, line_spacing=1.1)
    _source(slide, source_label, "Risk, implementation rule, and assumptions")

    for index, chunk in enumerate(research_chunks):
        page += 1
        title = "Recent Research Checks" if index == 0 else "Recent Research Checks (continued)"
        slide = _base_slide(prs, title, page, total, as_of)
        rows = [[item.ticker, _compact(item.development, 300), _compact(item.interpretation, 270), item.source_ids]
                for item in chunk]
        _table(slide, ["Ticker", "Recent material development", "Portfolio interpretation", "Sources"], rows,
               .75, 1.58, 11.83, 4.98, [.08, .42, .42, .08], header_size=9, body_size=8.1)
        _source(slide, source_label, "Research & Assumptions tab; source IDs cross-reference the audit")

    for index, chunk in enumerate(source_chunks):
        page += 1
        title = "Sources & Data Audit" if index == 0 else "Sources & Data Audit (continued)"
        slide = _base_slide(prs, title, page, total, as_of)
        rows = []
        for item in chunk:
            domain = urlsplit(item.url).netloc.removeprefix("www.") if item.url else ""
            source = item.source + (f" · {domain}" if domain else "")
            rows.append([
                item.source_id, _compact(item.item, 34), _compact(item.published, 18),
                _compact(source, 42), _compact(item.support, 120),
            ])
        _table(slide, ["ID", "Item", "Published", "Source", "What it supports"], rows,
               .75, 1.55, 11.83, 5.0, [.06, .18, .12, .24, .40], header_size=8.8, body_size=7.5)
        _source(slide, source_label, "URLs and evidence notes retained in the source workbook")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
