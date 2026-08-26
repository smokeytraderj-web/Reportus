"""Deterministic integrity checks before Reporticles enables finalization."""

from __future__ import annotations

import zipfile
from dataclasses import dataclass
from enum import StrEnum
from pathlib import Path


class QALevel(StrEnum):
    ERROR = "error"
    WARNING = "warning"


@dataclass(frozen=True, slots=True)
class QAIssue:
    level: QALevel
    code: str
    message: str


@dataclass(frozen=True, slots=True)
class OutputQAResult:
    approved: bool
    issues: tuple[QAIssue, ...] = ()
    page_or_sheet_count: int | None = None


class OutputInspector:
    """Verify a final artifact opens and has plausible structure."""

    def inspect(self, path: Path) -> OutputQAResult:
        issues: list[QAIssue] = []
        count: int | None = None
        if not path.is_file():
            return OutputQAResult(False, (QAIssue(QALevel.ERROR, "missing", "Output file is missing."),))
        if path.stat().st_size < 100:
            return OutputQAResult(False, (QAIssue(QALevel.ERROR, "too_small", "Output file is incomplete."),))
        try:
            extension = path.suffix.lower()
            if extension == ".pdf":
                count, issues = self._inspect_pdf(path)
            elif extension in {".xlsx", ".xlsm"}:
                count, issues = self._inspect_workbook(path)
            elif extension == ".pptx":
                count, issues = self._inspect_presentation(path)
            elif extension == ".docx":
                count, issues = self._inspect_docx(path)
            else:
                issues.append(QAIssue(QALevel.ERROR, "unsupported_output", "Output type is not supported."))
        except Exception:
            issues.append(QAIssue(QALevel.ERROR, "open_failed", "Output could not be opened for verification."))
        return OutputQAResult(
            approved=not any(issue.level is QALevel.ERROR for issue in issues),
            issues=tuple(issues),
            page_or_sheet_count=count,
        )

    def _inspect_pdf(self, path: Path) -> tuple[int, list[QAIssue]]:
        from pypdf import PdfReader

        reader = PdfReader(path)
        if reader.is_encrypted:
            return 0, [QAIssue(QALevel.ERROR, "encrypted", "Final PDF must not be encrypted.")]
        issues: list[QAIssue] = []
        for number, page in enumerate(reader.pages, start=1):
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            if width <= 0 or height <= 0:
                issues.append(QAIssue(QALevel.ERROR, "invalid_page", f"Page {number} has invalid dimensions."))
        if not reader.pages:
            issues.append(QAIssue(QALevel.ERROR, "no_pages", "Final PDF has no pages."))
        return len(reader.pages), issues

    def _inspect_workbook(self, path: Path) -> tuple[int, list[QAIssue]]:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=False, keep_links=False)
        try:
            visible = [sheet for sheet in workbook.worksheets if sheet.sheet_state == "visible"]
            issues = []
            if not visible:
                issues.append(QAIssue(QALevel.ERROR, "no_visible_sheets", "Workbook has no visible sheets."))
            if visible and not any(
                value is not None
                for sheet in visible
                for row in sheet.iter_rows(values_only=True)
                for value in row
            ):
                issues.append(QAIssue(QALevel.ERROR, "no_data", "Workbook contains no visible data."))
            return len(workbook.sheetnames), issues
        finally:
            workbook.close()

    def _inspect_presentation(self, path: Path) -> tuple[int, list[QAIssue]]:
        from pptx import Presentation

        presentation = Presentation(path)
        issues = []
        if not presentation.slides:
            issues.append(QAIssue(QALevel.ERROR, "no_slides", "Presentation contains no slides."))
        blank = sum(1 for slide in presentation.slides if not slide.shapes)
        if blank:
            issues.append(QAIssue(QALevel.WARNING, "blank_slides", f"Presentation contains {blank} blank slide(s)."))
        exact_prompts = {"slide number", "date", "footer", "paste chart here"}
        for slide_number, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape.text.strip()
                if getattr(shape, "is_placeholder", False) and not text:
                    issues.append(
                        QAIssue(
                            QALevel.ERROR,
                            "empty_placeholder",
                            f"Slide {slide_number} contains an unfilled placeholder.",
                        )
                    )
                normalized = text.casefold()
                if normalized in exact_prompts or normalized.startswith("click to add"):
                    issues.append(
                        QAIssue(
                            QALevel.ERROR,
                            "placeholder_text",
                            f"Slide {slide_number} contains unresolved placeholder text.",
                        )
                    )
        return len(presentation.slides), issues

    def _inspect_docx(self, path: Path) -> tuple[int, list[QAIssue]]:
        with zipfile.ZipFile(path) as archive:
            if "word/document.xml" not in archive.namelist():
                return 0, [QAIssue(QALevel.ERROR, "invalid_docx", "Word document structure is invalid.")]
        return 1, []
