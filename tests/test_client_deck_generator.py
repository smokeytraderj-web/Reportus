"""Synthetic end-to-end test for the Client Deck generator."""

import tempfile
import unittest
import shutil
from pathlib import Path

from generators.client_deck import AllocationRow, ClientDeckData, build_client_deck
from quality.output_qa import OutputInspector
from services.conversion import convert_pptx_to_pdf


class ClientDeckGeneratorTests(unittest.TestCase):
    def test_builds_structurally_valid_branded_deck(self) -> None:
        data = ClientDeckData(
            client_name="Sample Household",
            period="August 2026",
            as_of="As of August 25, 2026",
            allocation=(AllocationRow("Domestic Equity", 500000), AllocationRow("Fixed Income", 300000), AllocationRow("Cash", 200000)),
            risk_metrics={
                "Portfolio total": "$1,000,000", "Risk": "53",
                "Historical loss": "-$105,500", "Historical loss %": "-10.55%",
                "Historical gain": "+$189,700", "Historical gain %": "+18.97%",
                "Annual dividend": "2.26%", "Max drawdown": "-13.07%",
                "Annual range midpoint": "8.60%", "Expense ratio": "0.26%",
            },
            sector_performance={"Tech": .14, "Financials": .05, "Energy": .22},
            sector_portfolio={"Tech": 30.0, "Financials": 14.0, "Energy": 4.0},
            sector_benchmark={"Tech": 36.0, "Financials": 12.0, "Energy": 3.0},
            contributors=(("IVV", "S&P 500 ETF", "+13.37%", "+1.47%"),),
            detractors=(("TLT", "20+ Year Treasury ETF", "-2.61%", "-0.03%"),),
            earnings_years=("2024", "2025", "2026E", "2027E"),
            earnings_values=(250, 300, 380, 420),
            earnings_notes=("Earnings remain positive.", "Estimates are subject to revision."),
            optional_sections={"annuity": ("Current value: $200,000", "No contract fees reported.")},
            sources={"allocation": "Synthetic test data"},
        )
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "Sample_Client_Deck.pptx"

            build_client_deck(data, output)
            result = OutputInspector().inspect(output)

            self.assertTrue(result.approved, result.issues)
            self.assertEqual(result.page_or_sheet_count, 9)
            from pptx import Presentation
            risk_text = "\n".join(
                shape.text for shape in Presentation(output).slides[3].shapes
                if hasattr(shape, "text")
            )
            self.assertIn("95% HISTORICAL RANGE (6 MONTHS)", risk_text)
            self.assertIn("-$105,500", risk_text)
            self.assertIn("+$189,700", risk_text)
            if shutil.which("libreoffice") or shutil.which("soffice"):
                pdf = Path(directory) / "Sample_Client_Deck.pdf"
                convert_pptx_to_pdf(output, pdf)
                pdf_result = OutputInspector().inspect(pdf)
                self.assertTrue(pdf_result.approved, pdf_result.issues)
                self.assertEqual(pdf_result.page_or_sheet_count, 9)


if __name__ == "__main__":
    unittest.main()
