"""Tests for final-artifact integrity checks."""

import tempfile
import unittest
from pathlib import Path

from openpyxl import Workbook
from pypdf import PdfWriter
from pptx import Presentation

from quality.output_qa import OutputInspector


class OutputInspectorTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp_dir = tempfile.TemporaryDirectory()
        self.root = Path(self.temp_dir.name)
        self.inspector = OutputInspector()

    def tearDown(self) -> None:
        self.temp_dir.cleanup()

    def test_valid_pdf_is_approved(self) -> None:
        path = self.root / "report.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with path.open("wb") as stream:
            writer.write(stream)

        result = self.inspector.inspect(path)

        self.assertTrue(result.approved)
        self.assertEqual(result.page_or_sheet_count, 1)

    def test_valid_workbook_is_approved(self) -> None:
        path = self.root / "report.xlsx"
        workbook = Workbook()
        workbook.active.append(["Symbol", "Value"])
        workbook.active.append(["IVV", 100])
        workbook.save(path)

        result = self.inspector.inspect(path)

        self.assertTrue(result.approved)
        self.assertEqual(result.page_or_sheet_count, 1)

    def test_valid_presentation_is_approved(self) -> None:
        path = self.root / "report.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[0])
        presentation.save(path)

        result = self.inspector.inspect(path)

        self.assertTrue(result.approved)
        self.assertEqual(result.page_or_sheet_count, 1)

    def test_missing_output_is_rejected(self) -> None:
        result = self.inspector.inspect(self.root / "missing.pdf")

        self.assertFalse(result.approved)
        self.assertEqual(result.issues[0].code, "missing")


if __name__ == "__main__":
    unittest.main()
