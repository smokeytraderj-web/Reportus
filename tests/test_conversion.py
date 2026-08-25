"""Conversion adapter tests."""

import shutil
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from quality.output_qa import OutputInspector
from services.conversion import convert_pptx_to_pdf


@unittest.skipUnless(shutil.which("libreoffice") or shutil.which("soffice"), "Office converter unavailable")
class ConversionTests(unittest.TestCase):
    def test_pptx_converts_to_valid_pdf(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "source.pptx"
            destination = root / "result.pdf"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            slide.shapes.title.text = "Synthetic Reportus Test"
            presentation.save(source)

            convert_pptx_to_pdf(source, destination)

            self.assertTrue(OutputInspector().inspect(destination).approved)


if __name__ == "__main__":
    unittest.main()
