"""
Builds the blank GSWM "Client Deck" template (11 slides, no data).
Visual language matches Client Deck/layout/Mudry_Rinaldi_2026_Recommendations_Review_LAYOUT.pdf:
navy/gold, serif headlines, thin gold rule, clean white data pages.

This is a STRUCTURE-ONLY template. Every content slide (3-11) carries just a
title and a labeled placeholder zone -- no numbers, no charts, no tables of
real data. When a specific client's data is provided, populate a copy of
this deck (see output/build_<client>_deck.py for an example) using the
shared helpers in style.py so every client deck stays visually identical.

Section 9 was corrected from "RMS Report" to "RMD Report" (Required Minimum
Distribution) -- the firm's prep files for RMD tracking confirmed the report
type. Section 11 (Clear Spring Annuity) was added per firm request to cover
held-away/outside annuity assets that don't appear in the managed-account
allocation and sector reports.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from style import (
    NAVY, GOLD, WHITE, MUTED_ON_NAVY, BODY_GRAY, LIGHT_ROW, DARK_TEXT, TITLE_FONT, BODY_FONT,
    SLIDE_W, SLIDE_H,
    add_rect, add_gradient_bg, add_text, add_crest, add_footer, add_gold_rule,
    add_content_header, add_placeholder_zone, add_toc_slide,
)

OUT = r"C:\Users\Jack.Muirhead\OneDrive - Commonwealth Advisors\JJM Misc\Decky\Client Deck\GSWM_ClientDeck_Template.pptx"

CONTENT_SLIDES = [
    (3, "Overall Asset Allocation", "Asset allocation chart (by class) + allocation detail table"),
    (4, "Riskalyze Snapshot", "Riskalyze risk-score snapshot"),
    (5, "2026 Sector YTD Performance", "2026 sector year-to-date performance chart"),
    (6, "Equity Sector Exposure", "Sector allocation table (vs. S&P 500 TR USD) + relative diff bar chart"),
    (7, "Attribution Report", "Portfolio attribution report table"),
    (8, "S&P 500 Earnings Expectations", "Earnings growth chart, axis + YoY growth arrows"),
    (9, "RMD Report", "Required minimum distribution report"),
    (10, "Portfolio Summary \u2014 529 Accounts", "529 account portfolio summary table"),
    (11, "Clear Spring Annuity", "Clear Spring annuity summary"),
]
TOTAL = 11


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # --- Slide 1: Cover -----------------------------------------------
    s1 = prs.slides.add_slide(blank)
    add_gradient_bg(s1)
    add_crest(s1, SLIDE_W / 2, Inches(0.75))

    left = Inches(1.0)
    width = Inches(9.5)
    add_text(s1, left, Inches(3.55), width, Inches(0.3),
              "PRIVATE CLIENT REVIEW", 11, GOLD, font=BODY_FONT, bold=True, spacing=1.5)
    add_text(s1, left, Inches(3.9), width, Inches(0.85),
              "[Client / Household Name]", 40, WHITE, font=TITLE_FONT, bold=True)
    add_text(s1, left, Inches(4.62), width, Inches(0.5),
              "[Portfolio Review] \u00b7 [Period]", 22, WHITE, font=TITLE_FONT)
    add_gold_rule(s1, left, Inches(5.28), Inches(1.6), Pt(2.25))

    # --- Slide 2: Table of Contents ------------------------------------
    add_toc_slide(prs, blank, [(p, t) for p, t, _ in CONTENT_SLIDES], TOTAL)

    # --- Slides 3-11: content pages -------------------------------------
    for page_no, title, caption in CONTENT_SLIDES:
        s = prs.slides.add_slide(blank)
        add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
        add_content_header(s, title, page_no, TOTAL)
        add_placeholder_zone(s, caption)

    prs.save(OUT)
    print(f"Saved {OUT} with {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    build()
