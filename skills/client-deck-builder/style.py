"""
Shared visual system for the GSWM "Client Deck" style: navy/gold, serif
headlines, thin gold rule, clean white data pages. Matches
Client Deck/layout/Mudry_Rinaldi_2026_Recommendations_Review_LAYOUT.pdf.

Imported by build_template.py (blank master) and by per-client build scripts
under output/ so every client deck stays visually identical to the template.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_DARK = RGBColor(0x0A, 0x12, 0x24)
NAVY_MID = RGBColor(0x2B, 0x3D, 0x63)
GOLD = RGBColor(0xBF, 0xA0, 0x54)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# No neutral gray anywhere in the deck, including footers/page numbers: muted
# text on navy backgrounds uses warm gold-tan, muted text on white
# backgrounds is solid navy (same hue as titles, just small/italic so it
# still reads as secondary -- not a lighter/greyer shade of it).
MUTED_ON_NAVY = RGBColor(0xC9, 0xB8, 0x8A)
BODY_GRAY = RGBColor(0x1B, 0x2A, 0x4A)
FOOTER_GRAY = RGBColor(0x1B, 0x2A, 0x4A)
PLACEHOLDER_FILL = RGBColor(0xF7, 0xF7, 0xF9)
PLACEHOLDER_LINE = RGBColor(0xBF, 0xC4, 0xCE)
PLACEHOLDER_TEXT = RGBColor(0x1B, 0x2A, 0x4A)
LIGHT_ROW = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
GREEN = RGBColor(0x0B, 0x7A, 0x32)
RED = RGBColor(0xB0, 0x00, 0x20)
ACCENT_BLUE = RGBColor(0x2F, 0x74, 0xC4)
TABLE_HEADER_LIGHT = RGBColor(0xEC, 0xEF, 0xF4)

# chart palette: navy family + gold, used across every native chart so every
# slide's chart reads as one system rather than PowerPoint defaults. No
# neutral gray in the mix -- the old 6th color read as "grey" on screen.
CHART_COLORS = [NAVY, GOLD, RGBColor(0x6E, 0x82, 0xA8), RGBColor(0xD9, 0xC4, 0x8A),
                 RGBColor(0x40, 0x55, 0x7A), RGBColor(0x8F, 0x6B, 0x2E)]

TITLE_FONT = "Garamond"
BODY_FONT = "Times New Roman"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_letter_spacing(run, pts):
    try:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(pts * 100)))
    except Exception:
        pass


def no_line(shape):
    shape.line.fill.background()


def no_shadow(shape):
    shape.shadow.inherit = False


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_pt=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    no_shadow(shp)
    if fill_rgb is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        no_line(shp)
    else:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(line_pt or 0.75)
    return shp


def add_gradient_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    no_shadow(bg)
    no_line(bg)
    bg.fill.gradient()
    stops = bg.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = NAVY_MID
    stops[1].position = 1.0
    stops[1].color.rgb = NAVY_DARK
    try:
        bg.fill.gradient_angle = 45.0
    except Exception:
        pass
    return bg


def add_text(slide, left, top, width, height, text, size, color, font=BODY_FONT,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             spacing=None, line_spacing=None, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if spacing:
        set_letter_spacing(r, spacing)
    return box


def add_bullets(slide, left, top, width, height, items, size=14, color=DARK_TEXT,
                 font=BODY_FONT, space_after=10, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def add_crest(slide, center_x, top):
    outer_d = Inches(1.5)
    inner_d = Inches(1.22)
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - outer_d / 2, top, outer_d, outer_d)
    outer.fill.background()
    outer.line.color.rgb = GOLD
    outer.line.width = Pt(1.25)
    no_shadow(outer)
    inner_top = top + (outer_d - inner_d) / 2
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - inner_d / 2, inner_top, inner_d, inner_d)
    inner.fill.background()
    inner.line.color.rgb = GOLD
    inner.line.width = Pt(0.75)
    no_shadow(inner)
    add_text(slide, center_x - Inches(0.75), top + outer_d / 2 - Inches(0.28), Inches(1.5), Inches(0.56),
              "GS", 28, GOLD, font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, center_x - Inches(2.25), top + outer_d + Inches(0.12), Inches(4.5), Inches(0.3),
              "GOTTFRIED & SOMBERG WEALTH MANAGEMENT", 9, GOLD, font=BODY_FONT,
              align=PP_ALIGN.CENTER, spacing=1.5)


def add_footer(slide, page_no, total):
    add_text(slide, Inches(0.75), Inches(7.14), Inches(5.0), Inches(0.25),
              "Gottfried & Somberg Wealth Management", 9, FOOTER_GRAY, font=BODY_FONT)
    add_text(slide, Inches(9.583), Inches(7.14), Inches(3.0), Inches(0.25),
              f"Page {page_no} of {total}", 9, FOOTER_GRAY, font=BODY_FONT, align=PP_ALIGN.RIGHT)


def add_gold_rule(slide, left, top, width, height=Pt(2.25)):
    return add_rect(slide, left, top, width, height, fill_rgb=GOLD)


def add_content_header(slide, title, page_no, total, meta_placeholder="[As of Date]", title_size=30):
    add_text(slide, Inches(0.75), Inches(0.55), Inches(8.5), Inches(0.6),
              title, title_size, NAVY, font=TITLE_FONT, bold=True)
    add_text(slide, Inches(8.583), Inches(0.62), Inches(4.0), Inches(0.4),
              meta_placeholder, 12, BODY_GRAY, font=BODY_FONT, italic=True, align=PP_ALIGN.RIGHT)
    add_gold_rule(slide, Inches(0.75), Inches(1.28), Inches(11.833))
    add_footer(slide, page_no, total)


def add_placeholder_zone(slide, caption, left=Inches(0.75), top=Inches(1.62),
                           width=Inches(11.833), height=Inches(5.15)):
    ph = add_rect(slide, left, top, width, height, fill_rgb=PLACEHOLDER_FILL, line_rgb=PLACEHOLDER_LINE, line_pt=1.0)
    try:
        ln = ph.line._get_or_add_ln()
        prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(prstDash)
    except Exception:
        pass
    tf = ph.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = caption
    r.font.name = BODY_FONT
    r.font.size = Pt(15)
    r.font.italic = True
    r.font.color.rgb = PLACEHOLDER_TEXT
    return ph


def style_table(table, headers, rows, col_weights=None, total_width=Inches(11.833),
                 header_size=12, body_size=11.5, align_from_col=None, header_align_from_col=None):
    """Flat firm-style table: navy header row, white/light-gray zebra body,
    no gridlines. align_from_col: 0-indexed column at/after which cells are
    right-aligned (for numeric columns)."""
    n_cols = len(headers)
    if col_weights:
        for c, w in enumerate(col_weights):
            table.columns[c].width = int(total_width * w)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
        cell.margin_left = Pt(8)
        cell.margin_right = Pt(8)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if (header_align_from_col is not None and c >= header_align_from_col) else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = h
        r.font.name = BODY_FONT
        r.font.bold = True
        r.font.size = Pt(header_size)
        r.font.color.rgb = WHITE

    for i, row in enumerate(rows):
        bg = WHITE if i % 2 == 0 else LIGHT_ROW
        for c in range(n_cols):
            cell = table.cell(i + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if (align_from_col is not None and c >= align_from_col) else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(row[c])
            r.font.name = BODY_FONT
            r.font.size = Pt(body_size)
            r.font.bold = bool(row[c]) and str(row[c]).endswith("Total") if False else False
            r.font.color.rgb = DARK_TEXT

    tbl_el = table._tbl
    tblPr = tbl_el.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')


def add_toc_slide(prs, blank, sections, total):
    """Full-bleed navy TOC: numbered list (not a table), matches the cover's
    dark treatment so cover -> contents reads as one continuous opening
    instead of dark-then-suddenly-white. No subtitle line -- section list
    starts right under the title."""
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_text(s, Inches(0.9), Inches(0.65), Inches(8), Inches(0.3),
              "PORTFOLIO REVIEW OVERVIEW", 11, GOLD, font=BODY_FONT, bold=True, spacing=1.5)
    add_text(s, Inches(0.9), Inches(0.97), Inches(8), Inches(0.75),
              "Contents", 34, WHITE, font=TITLE_FONT, bold=True)

    list_top = Inches(1.95)
    row_h = Inches(0.52)
    num_w = Inches(0.7)
    for i, (page_no, title) in enumerate(sections):
        row_top = list_top + row_h * i
        add_text(s, Inches(0.9), row_top, num_w, row_h,
                  f"{i + 1:02d}", 16, GOLD, font=TITLE_FONT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.75), row_top, Inches(8.5), row_h,
                  title, 15.5, WHITE, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(10.4), row_top, Inches(1.7), row_h,
                  f"Page {page_no}", 12, MUTED_ON_NAVY, font=BODY_FONT, align=PP_ALIGN.RIGHT,
                  anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.9), Inches(7.14), Inches(6.0), Inches(0.25),
              "Gottfried & Somberg Wealth Management", 9, MUTED_ON_NAVY, font=BODY_FONT)
    add_text(s, Inches(9.583), Inches(7.14), Inches(2.85), Inches(0.25),
              f"Page 2 of {total}", 9, MUTED_ON_NAVY, font=BODY_FONT, align=PP_ALIGN.RIGHT)
    return s


def add_diff_bars_manual(slide, left, top, width, height, categories, values,
                           pos_color=NAVY, neg_color=GOLD):
    """Same over/underweight-vs-zero visual as add_diff_bar_chart, but drawn
    as plain rectangle shapes instead of a native chart. PowerPoint's chart
    engine has a reproducible bug in this environment where negative-value
    columns export with no fill (confirmed on the underlying XML being
    correct -- dPt overrides, two-series overlap, and a forced chart-type
    refresh via COM all failed to fix it), so this sidesteps the chart
    engine entirely for this one visual. Plain shapes render fill correctly
    everywhere else in this deck."""
    n = len(categories)
    slot_w = width / n
    bar_w = slot_w * 0.55
    label_row_h = Inches(0.26)
    # Fixed (not proportional) headroom reserved above the tallest positive
    # bar and below the tallest negative bar for their value-labels, so a
    # small value span can't starve the label of room the way a percentage
    # padding would.
    end_margin = Inches(0.35)

    pos_max = max([v for v in values if v > 0], default=0)
    neg_max = abs(min([v for v in values if v < 0], default=0))
    value_span = (pos_max + neg_max) or 1
    bar_area_h = height - label_row_h - 2 * end_margin
    ppu = bar_area_h / value_span  # EMU per data unit
    zero_y = top + end_margin + int(pos_max * ppu)

    add_rect(slide, left, zero_y, width, Pt(1.25), fill_rgb=NAVY)

    for i, (cat, val) in enumerate(zip(categories, values)):
        cx = left + slot_w * i + (slot_w - bar_w) / 2
        color = pos_color if val >= 0 else neg_color
        if val >= 0:
            bar_h = int(val * ppu)
            bar_top = zero_y - bar_h
        else:
            bar_h = int(abs(val) * ppu)
            bar_top = zero_y + label_row_h
        if bar_h > 0:
            add_rect(slide, cx, bar_top, bar_w, bar_h, fill_rgb=color)

        label_text = f"{val:+.2f}"
        if val >= 0:
            add_text(slide, left + slot_w * i, bar_top - Inches(0.28), slot_w, Inches(0.24),
                      label_text, 9.5, NAVY, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
        else:
            add_text(slide, left + slot_w * i, bar_top + bar_h + Inches(0.06), slot_w, Inches(0.24),
                      label_text, 9.5, NAVY, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)

        add_text(slide, left + slot_w * i, zero_y + Inches(0.02), slot_w, label_row_h,
                  cat, 9.5, NAVY, font=BODY_FONT, align=PP_ALIGN.CENTER)


def add_donut_chart(slide, left, top, width, height, categories, values, colors=None, hole_size=65):
    # add_chart, like add_connector, serializes its xfrm offsets without an
    # int cast -- a float EMU (e.g. from a `/` division upstream) writes a
    # literal trailing ".0" into the XML, an invalid ST_Coordinate that
    # PowerPoint reports as file corruption on open.
    left, top, width, height = int(left), int(top), int(width), int(height)
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Allocation", values)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left, top, width, height, data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.number_format = '0.0%'
    dls.number_format_is_linked = False
    dls.font.size = Pt(11)
    dls.font.name = BODY_FONT
    dls.font.color.rgb = NAVY
    try:
        chart.plots[0].donut_hole_size = hole_size
    except Exception:
        pass
    series = plot.series[0]
    palette = colors or CHART_COLORS
    for i, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = palette[i % len(palette)]
        point.format.line.color.rgb = WHITE
        point.format.line.width = Pt(1.5)
    return gframe


def add_bar_chart(slide, left, top, width, height, categories, values, color=NAVY,
                    number_format='0.0%', title=None, data_label_color=None):
    left, top, width, height = int(left), int(top), int(width), int(height)
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Series 1", values)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.number_format = number_format
    dls.number_format_is_linked = False
    dls.font.size = Pt(10.5)
    dls.font.bold = True
    dls.font.name = BODY_FONT
    dls.font.color.rgb = data_label_color or NAVY
    dls.position = XL_LABEL_POSITION.OUTSIDE_END
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.fill.background()

    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(10)
    cat_ax.tick_labels.font.name = BODY_FONT
    cat_ax.format.line.color.rgb = RGBColor(0xBF, 0xC4, 0xCE)
    cat_ax.major_tick_mark = XL_TICK_MARK.NONE

    val_ax = chart.value_axis
    val_ax.visible = False
    val_ax.has_major_gridlines = False
    return gframe


def add_axis_bar_chart(slide, left, top, width, height, categories, values, colors,
                         value_min=0, value_max=None, major_unit=None, number_format='0'):
    """Column chart with a visible value axis (gridlines + tick labels), used
    where the reference chart shows the reader the scale directly rather than
    via OUTSIDE_END data labels. Per-point fill via dPt overrides -- safe
    here because every value is positive (the negative-dPt export bug noted
    on add_diff_bars_manual only hits negative-value points)."""
    left, top, width, height = int(left), int(top), int(width), int(height)
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Series 1", values)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 45
    plot.has_data_labels = False

    series = plot.series[0]
    series.format.line.fill.background()
    for i, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[i % len(colors)]
        point.format.line.fill.background()

    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(11)
    cat_ax.tick_labels.font.bold = True
    cat_ax.tick_labels.font.name = BODY_FONT
    cat_ax.format.line.color.rgb = RGBColor(0x8A, 0x90, 0x9C)
    cat_ax.major_tick_mark = XL_TICK_MARK.NONE

    val_ax = chart.value_axis
    val_ax.visible = True
    val_ax.minimum_scale = value_min
    if value_max is not None:
        val_ax.maximum_scale = value_max
    if major_unit is not None:
        val_ax.major_unit = major_unit
    val_ax.has_major_gridlines = True
    val_ax.major_gridlines.format.line.color.rgb = RGBColor(0xE3, 0xE6, 0xEC)
    val_ax.major_gridlines.format.line.width = Pt(0.75)
    val_ax.format.line.color.rgb = RGBColor(0x8A, 0x90, 0x9C)
    val_ax.major_tick_mark = XL_TICK_MARK.NONE
    val_ax.tick_labels.font.size = Pt(10)
    val_ax.tick_labels.font.name = BODY_FONT
    val_ax.tick_labels.font.color.rgb = NAVY
    val_ax.tick_labels.number_format = number_format
    val_ax.tick_labels.number_format_is_linked = False
    return gframe


def add_growth_arrow(slide, x1, y1, x2, y2, label, label_width=Inches(0.9)):
    """Diagonal green growth arrow between two bar tops, with a bold percent
    label centered above its midpoint -- matches the reference chart's YoY
    callouts. Coordinates are coerced to int EMU: unlike add_rect/add_shape,
    add_connector serializes its xfrm offsets without an int cast, so a
    float EMU value (e.g. from a `/` division upstream) gets written to the
    XML with a literal trailing ".0" -- an invalid ST_Coordinate that
    PowerPoint reports as file corruption on open."""
    x1, y1, x2, y2 = int(x1), int(y1), int(x2), int(y2)
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = MSO_CONNECTOR.STRAIGHT
    conn.line.color.rgb = GREEN
    conn.line.width = Pt(2.25)
    no_shadow(conn)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    mid_x = (x1 + x2) / 2
    mid_y = min(y1, y2)
    add_text(slide, mid_x - label_width / 2, mid_y - Inches(0.32), label_width, Inches(0.26),
              label, 13, GREEN, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    return conn


def add_light_comparison_table(slide, left, top, total_width, header_h, row_h, label_w,
                                 categories, row_labels, rows, divider_before=None):
    """Light-gray-header comparison table (vs. the navy-block header used
    elsewhere in this deck) -- one header row of category names, then one
    row per series in `rows` (aligned to `row_labels`). `divider_before` draws
    a thin vertical rule before that category index, e.g. to set off
    aggregate "Developed/Emerging" columns from the individual regions."""
    n = len(categories)
    n_rows = 1 + len(row_labels)
    tbl_shape = slide.shapes.add_table(n_rows, n + 1, left, top,
                                         total_width, header_h + row_h * len(row_labels))
    table = tbl_shape.table
    table.columns[0].width = label_w
    col_w = int((total_width - label_w) / n)
    for c in range(n):
        table.columns[c + 1].width = col_w
    table.rows[0].height = header_h
    for r in range(len(row_labels)):
        table.rows[r + 1].height = row_h

    def set_cell(r, c, text, fill, color, bold=True, size=9.5, align=PP_ALIGN.CENTER):
        cell = table.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(3)
        cell.margin_right = Pt(3)
        cell.margin_top = Pt(2)
        cell.margin_bottom = Pt(2)
        tf = cell.text_frame
        tf.word_wrap = True
        lines = text.split("\n")
        for li, line in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.alignment = align
            r_ = p.add_run()
            r_.text = line
            r_.font.name = BODY_FONT
            r_.font.bold = bold
            r_.font.size = Pt(size)
            r_.font.color.rgb = color

    set_cell(0, 0, "", TABLE_HEADER_LIGHT, NAVY)
    for c, cat in enumerate(categories):
        set_cell(0, c + 1, cat, TABLE_HEADER_LIGHT, NAVY, size=8.5)
    for ri, label in enumerate(row_labels):
        bg = WHITE if ri % 2 == 0 else LIGHT_ROW
        set_cell(ri + 1, 0, label, bg, NAVY, align=PP_ALIGN.LEFT, size=9.5)
        for c, val in enumerate(rows[ri]):
            set_cell(ri + 1, c + 1, f"{val:.2f}", bg, DARK_TEXT, bold=False, size=9.5)

    tbl_el = table._tbl
    tblPr = tbl_el.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')

    if divider_before is not None:
        div_x = left + label_w + col_w * divider_before
        add_rect(slide, div_x, top, Pt(1.25), header_h + row_h * len(row_labels), fill_rgb=NAVY)

    return tbl_shape, col_w


def add_diff_bars_single_tone(slide, left, top, width, height, categories, values,
                                color=ACCENT_BLUE, legend_label=None):
    """Zero-baseline diff bars in one color regardless of sign (matches the
    'Equity regional exposure' reference: a single blue series, not a
    pos/neg two-color split), with a fixed bottom row of category labels
    (Excel's axis-labels-at-Low placement) so a single oversized outlier bar
    can't push every label up into the plot area, plus an optional
    dot+text legend above the chart."""
    n = len(categories)
    slot_w = width / n
    bar_w = slot_w * 0.42
    label_row_h = Inches(0.42)
    val_label_h = Inches(0.24)
    end_margin = Inches(0.08)

    if legend_label:
        dot_d = Inches(0.09)
        dot_x = left + width - Inches(2.6)
        dot_y = top - Inches(0.30)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, dot_y, dot_d, dot_d)
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        no_line(dot)
        no_shadow(dot)
        add_text(slide, dot_x + dot_d + Inches(0.08), dot_y - Inches(0.06), Inches(3.2), Inches(0.22),
                  legend_label, 9.5, NAVY, font=BODY_FONT, bold=True)

    pos_max = max([v for v in values if v > 0], default=0)
    neg_max = abs(min([v for v in values if v < 0], default=0))
    value_span = (pos_max + neg_max) or 1
    # Reserve val_label_h on BOTH ends of the bar area -- a positive bar's
    # value label sits above it (needs top headroom) and a negative bar's
    # value label sits below it, right before the fixed category-label row
    # (needs bottom headroom too), or the biggest negative outlier's label
    # collides with the category names under it.
    bar_area_h = height - label_row_h - 2 * val_label_h - end_margin
    ppu = bar_area_h / value_span
    zero_y = top + val_label_h + int(pos_max * ppu)

    add_rect(slide, left, zero_y, width, Pt(1.0), fill_rgb=RGBColor(0x40, 0x40, 0x40))

    label_top = top + height - label_row_h
    for i, (cat, val) in enumerate(zip(categories, values)):
        cx = left + slot_w * i + (slot_w - bar_w) / 2
        if val >= 0:
            bar_h = max(int(val * ppu), 1)
            bar_top = zero_y - bar_h
            lbl_top = bar_top - Inches(0.27)
        else:
            bar_h = max(int(abs(val) * ppu), 1)
            bar_top = zero_y
            lbl_top = bar_top + bar_h + Inches(0.03)
        add_rect(slide, cx, bar_top, bar_w, bar_h, fill_rgb=color)

        add_text(slide, left + slot_w * i, lbl_top, slot_w, Inches(0.22),
                  f"{val:+.2f}", 8, DARK_TEXT, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)

        box = slide.shapes.add_textbox(left + slot_w * i - Inches(0.04), label_top,
                                         slot_w + Inches(0.08), label_row_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        lines = cat.split("\n")
        for li, line in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.0
            r_ = p.add_run()
            r_.text = line
            r_.font.name = BODY_FONT
            r_.font.size = Pt(7.5)
            r_.font.color.rgb = NAVY


def add_diff_bar_chart(slide, left, top, width, height, categories, values, pos_color=NAVY,
                         neg_color=GOLD, number_format='+0.00;-0.00'):
    """Bars rising/falling from a zero baseline (over/underweight vs.
    benchmark, YoY change, etc.) -- category axis sits at zero so it doubles
    as the baseline, matching the reference JPM 'relative to benchmark'
    chart. Colored by sign using two fully-overlapped series (one positive,
    one negative, blank everywhere the other has data) rather than per-point
    dPt overrides on a single series -- PowerPoint's chart cache does not
    reliably re-render single-series dPt fills through headless COM export,
    so it exports as an outline with no fill; two plain-colored series avoid
    that entirely."""
    left, top, width, height = int(left), int(top), int(width), int(height)
    pos_vals = [v if v >= 0 else None for v in values]
    neg_vals = [v if v < 0 else None for v in values]
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Positive", pos_vals)
    data.add_series("Negative", neg_vals)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 50
    plot.overlap = 100
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.number_format = number_format
    dls.number_format_is_linked = False
    dls.font.size = Pt(9.5)
    dls.font.bold = True
    dls.font.name = BODY_FONT
    dls.font.color.rgb = NAVY

    pos_series, neg_series = plot.series[0], plot.series[1]
    for series, color in ((pos_series, pos_color), (neg_series, neg_color)):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        series.format.line.fill.background()

    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(9.5)
    cat_ax.tick_labels.font.name = BODY_FONT
    cat_ax.format.line.color.rgb = NAVY
    cat_ax.format.line.width = Pt(1)
    cat_ax.major_tick_mark = XL_TICK_MARK.NONE

    val_ax = chart.value_axis
    val_ax.has_major_gridlines = False
    val_ax.format.line.fill.background()
    val_ax.tick_labels.font.size = Pt(1)
    return gframe
