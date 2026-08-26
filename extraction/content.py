"""Extract bounded, provenance-tagged text from privacy-approved uploads."""

from __future__ import annotations

import csv
import json
import re
import shutil
import subprocess
import tempfile
from pathlib import Path

from PIL import Image, ImageFilter, ImageOps

from providers.base import SourceFragment


class ExtractionError(RuntimeError):
    pass


_SIGNED_AMOUNT = re.compile(r"[+-]\s*\$\s*\d[\d,\s]*\d")
_DECIMAL_PERCENT = re.compile(r"\+?\d+\.\d+%")


def _normalize_riskalyze_range(text: str) -> str:
    """Add deterministic labels to the four values in Riskalyze's range card."""

    amounts = [re.sub(r"\s+", "", value) for value in _SIGNED_AMOUNT.findall(text)]
    percentages = _DECIMAL_PERCENT.findall(text)
    loss = next((value for value in amounts if value.startswith("-")), "")
    gain = next((value for value in amounts if value.startswith("+")), "")
    if not loss or not gain or len(percentages) < 2:
        return text
    loss_percent = percentages[0].lstrip("+")
    gain_percent = percentages[1].lstrip("+")
    labels = (
        f"Historical loss: {loss}",
        f"Historical loss %: -{loss_percent}",
        f"Historical gain: {gain}",
        f"Historical gain %: +{gain_percent}",
    )
    return text.rstrip() + "\n\n" + "\n".join(labels) + "\n"


class ContentExtractor:
    def __init__(self, *, max_characters_per_file: int = 200_000, max_total_characters: int = 500_000):
        self.max_characters_per_file = max_characters_per_file
        self.max_total_characters = max_total_characters

    def extract(self, paths: tuple[Path, ...] | list[Path]) -> tuple[SourceFragment, ...]:
        fragments: list[SourceFragment] = []
        total = 0
        for path in paths:
            extracted = self._extract_file(path)
            file_total = sum(len(item.text) for item in extracted)
            if file_total > self.max_characters_per_file:
                raise ExtractionError(f"{path.name} is too large for one synthesis request.")
            total += file_total
            if total > self.max_total_characters:
                raise ExtractionError("The selected sources are too large for one synthesis request.")
            fragments.extend(extracted)
        if not fragments:
            raise ExtractionError("No readable source content was found.")
        return tuple(fragments)

    def _extract_file(self, path: Path) -> list[SourceFragment]:
        extension = path.suffix.lower()
        if extension in {".txt", ".md", ".json", ".xml", ".html", ".htm"}:
            text = path.read_text(encoding="utf-8-sig", errors="strict")
            if extension == ".json":
                text = json.dumps(json.loads(text), indent=2, ensure_ascii=False)
            return [SourceFragment(path.name, "document", text)]
        if extension in {".csv", ".tsv"}:
            delimiter = "\t" if extension == ".tsv" else ","
            with path.open("r", encoding="utf-8-sig", newline="") as stream:
                rows = list(csv.reader(stream, delimiter=delimiter))
            return [
                SourceFragment(path.name, f"row {number}", " | ".join(str(value) for value in row))
                for number, row in enumerate(rows, start=1)
            ]
        if extension in {".xlsx", ".xlsm"}:
            return self._extract_workbook(path)
        if extension == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(path)
            return [
                SourceFragment(path.name, f"page {number}", page.extract_text() or "")
                for number, page in enumerate(reader.pages, start=1)
            ]
        if extension == ".docx":
            from docx import Document

            document = Document(path)
            items = [SourceFragment(path.name, f"paragraph {number}", paragraph.text) for number, paragraph in enumerate(document.paragraphs, start=1) if paragraph.text.strip()]
            for table_number, table in enumerate(document.tables, start=1):
                for row_number, row in enumerate(table.rows, start=1):
                    items.append(SourceFragment(path.name, f"table {table_number}, row {row_number}", " | ".join(cell.text for cell in row.cells)))
            return items
        if extension == ".pptx":
            from pptx import Presentation

            presentation = Presentation(path)
            return [
                SourceFragment(path.name, f"slide {number}", "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()))
                for number, slide in enumerate(presentation.slides, start=1)
            ]
        if extension in {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}:
            text = self._ocr_riskalyze(path) if path.stem.startswith("riskalyze_analytics_") else self._ocr(path)
            return [SourceFragment(path.name, "image OCR", text)]
        raise ExtractionError(f"{path.name} is not supported for local text extraction.")

    def _extract_workbook(self, path: Path) -> list[SourceFragment]:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=True, keep_links=False)
        fragments: list[SourceFragment] = []
        try:
            for sheet in workbook.worksheets:
                if sheet.sheet_state != "visible":
                    continue
                for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    if not any(value is not None for value in row):
                        continue
                    text = " | ".join("" if value is None else str(value) for value in row)
                    fragments.append(SourceFragment(path.name, f"{sheet.title} row {row_number}", text))
        finally:
            workbook.close()
        return fragments

    def _ocr(self, path: Path) -> str:
        executable = shutil.which("tesseract")
        if executable is None:
            raise ExtractionError("Image OCR is unavailable.")
        with tempfile.TemporaryDirectory(prefix="reportus-extract-") as temporary:
            output = Path(temporary) / "ocr"
            completed = subprocess.run(
                [executable, str(path), str(output)], capture_output=True, text=True,
                timeout=90, check=False,
            )
            text_path = output.with_suffix(".txt")
            if completed.returncode != 0 or not text_path.is_file():
                raise ExtractionError("Image OCR failed.")
            return text_path.read_text(encoding="utf-8", errors="replace")

    def _ocr_riskalyze(self, path: Path) -> str:
        """OCR the safe analytics crop plus an enlarged historical-range card."""

        executable = shutil.which("tesseract")
        if executable is None:
            raise ExtractionError("Image OCR is unavailable.")
        with tempfile.TemporaryDirectory(prefix="reportus-riskalyze-ocr-") as temporary:
            root = Path(temporary)
            with Image.open(path) as raw:
                image = ImageOps.autocontrast(ImageOps.grayscale(raw)).filter(ImageFilter.SHARPEN)
                full_path = root / "analytics.png"
                image.save(full_path)
                range_panel = image.crop((0, round(image.height * .38), image.width, round(image.height * .58)))
                range_panel = range_panel.resize(
                    (round(range_panel.width * 2.5), round(range_panel.height * 2.5)),
                    Image.Resampling.LANCZOS,
                )
                range_path = root / "historical-range.png"
                range_panel.save(range_path)

            outputs = []
            for image_path, psm in ((full_path, "4"), (range_path, "6")):
                completed = subprocess.run(
                    [executable, str(image_path), "stdout", "--psm", psm],
                    capture_output=True, text=True, timeout=90, check=False,
                )
                if completed.returncode != 0:
                    raise ExtractionError("Riskalyze image OCR failed.")
                outputs.append(completed.stdout)
            return outputs[0].rstrip() + "\n\n" + _normalize_riskalyze_range(outputs[1])
